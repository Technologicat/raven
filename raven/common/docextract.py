"""Extract indexable plaintext from document files (plain text, PDF, office documents, and HTML).

This is the single source of truth for "given a document file, give me its text" across Raven. The RAG
document-database ingester (`raven.librarian.hybridir`), the chat document attachment feature, and the
`raven-pdf2bib` tool all route through here, so there is exactly one backend per format and one
supported-format list. Adding a format here adds it to every one of those surfaces at once.

Contract of `extract_text`:

  - raises `FileNotFoundError` if the path does not exist;
  - raises `DocumentExtractionError` if the file exists but cannot be parsed as its format (a corrupt or
    encrypted PDF, a text file that is not valid UTF-8) — the underlying cause is chained;
  - returns `None` if the file parses cleanly but yields no text (a scanned PDF with no text layer, a
    whitespace-only text file);
  - otherwise returns the extracted text.

The raise-vs-`None` split is deliberate: an *error situation* (missing or unreadable) is a different thing from
an *empty but valid* document. The rule is *parse failure → raise; parsed-but-empty → `None`*. Callers apply
their own policy on top of it — a background batch ingester catches the exceptions and skips the offending file
so one bad document does not abort the batch, whereas an interactive attach site lets the exception surface so
it can tell the user *why* their file could not be read.

PDF text extraction handles born-digital PDFs (a real text layer). A scanned/image-only PDF has no text to
extract and comes back as `None`; OCR for those is a separate, later concern.

HTML goes through the same readability extraction the web-fetch tool uses, so a page saved to disk reads the way
the same page fetched live would. What that misses is a page whose text is produced by *running* it: the bare
shell of a JS-rendered site, and equally a self-contained single-file app that carries its data inline as a
script literal and builds the DOM at load. In the second case the content genuinely is in the file — it is just
in a `<script>` element, which readability extraction ignores, and reading it would mean deciding where inline
data ends and a minified bundle begins. Neither is attempted here; both extract as empty, the way a scanned PDF
does.

Every format here is read for its *text layer* only. Whatever a document says through pictures — a figure, a
photograph, a typeset equation rendered as an image — is not recovered, and a file whose content is entirely
such material extracts as empty. Legacy binary office formats (`.doc`, `.ppt`) are not supported: reading them
means an external converter process, and PDF handling was deliberately moved off one of those.
"""

__all__ = ["DocumentExtractionError",
           "supported_extensions",
           "is_supported",
           "repair_surrogates",
           "extract_text", "extract_text_from_bytes",
           "Extractor",
           "PLAINTEXT", "ALL_FORMATS"]

import io
import logging
import pathlib
from typing import Any, BinaryIO, Callable, Iterable, Iterator

import docx
import docx.table
import odf.namespaces
import odf.opendocument
import odf.teletype
import pptx
import pptx.enum.shapes
import pypdf

from . import text as common_text

logger = logging.getLogger(__name__)


class DocumentExtractionError(Exception):
    """A document file exists but its text could not be extracted (corrupt/encrypted PDF, wrong text encoding)."""


# Extensions read verbatim as UTF-8 plain text. Beyond `.txt` these are all markup a language model reads
# perfectly well as-is, so the list is really a list of who turns up carrying which format:
#
#   - `.md`, and `.rst` as its cousin — the latter is the scientific Python stack's house markup, so it arrives
#     with anyone whose reading is open-source documentation;
#   - `.bib`, and `.tex` alongside it, because someone who has the bibliography almost always has the papers
#     that cite it as well;
#   - `.org`, for those who count the Church of Emacs among their affiliations.
#
# Keep in sync with the formats the librarian offers for ingestion (`raven.librarian.config.llm_docs_exts`) —
# this tuple is the extractor's own capability, that config is the user-facing enable list.
_PLAINTEXT_EXTS = (".txt", ".md", ".rst", ".org", ".bib", ".tex")


# Every extractor below takes an open binary stream and a `label` naming where it came from — a path for a
# file, the original filename for an attachment held in memory. The label is used only in messages and log
# lines; nothing dispatches on it. This is what lets one backend serve both entry points, and it costs
# nothing, because every library underneath reads a stream (verified: pypdf, python-docx, python-pptx,
# odfpy and trafilatura all accept one).
def _extract_plaintext(source: BinaryIO, label: str) -> str:
    try:
        return source.read().decode("utf-8")
    except UnicodeDecodeError as exc:
        raise DocumentExtractionError(f"extract_text: '{label}' is not valid UTF-8 text.") from exc


def _extract_pdf(source: BinaryIO, label: str) -> str:
    # Force the whole page list up front so encryption/corruption errors (which pypdf raises lazily on page
    # access) surface here, as one clean `DocumentExtractionError`, rather than mid-iteration below.
    try:
        reader = pypdf.PdfReader(source)
        pages = list(reader.pages)
    except Exception as exc:  # noqa: BLE001 -- pypdf raises many types on malformed/encrypted input; normalize them
        raise DocumentExtractionError(f"extract_text: '{label}' could not be read as a PDF: "
                                      f"{type(exc).__name__}: {exc}") from exc
    chunks = []
    for page_number, page in enumerate(pages):
        try:
            chunks.append(page.extract_text())
        except Exception as exc:  # noqa: BLE001 -- one unreadable page must not lose the rest of the document
            logger.warning(f"extract_text: '{label}': skipping unreadable page {page_number}: "
                           f"{type(exc).__name__}: {exc}")
    return "\n".join(chunks)


# Cells of one table row are joined by a tab rather than a newline, so that a row survives as a row. Retrieval
# chunks on a sliding window of text, so a table flattened one-cell-per-line loses the association between a
# label and its value — which for a table is most of what it says.
_TABLE_CELL_SEPARATOR = "\t"


def _docx_content_text(container: Any) -> str:
    """Text of a `python-docx` block container — a whole document, or one table cell — in document order.

    Recursive, because a table cell is itself a block container: a table nested in a cell comes out in place.
    """
    chunks = []
    for block in container.iter_inner_content():
        if isinstance(block, docx.table.Table):
            for row in block.rows:
                chunks.append(_TABLE_CELL_SEPARATOR.join(_docx_content_text(cell) for cell in row.cells))
        else:  # a paragraph
            chunks.append(block.text)
    return "\n".join(chunks)


def _extract_docx(source: BinaryIO, label: str) -> str:
    # Headers, footers and comments are deliberately not collected: they repeat on every page or annotate rather
    # than state, so folding them into the body text mostly adds boilerplate for a retrieval index to match on.
    try:
        return _docx_content_text(docx.Document(source))
    except Exception as exc:  # noqa: BLE001 -- python-docx surfaces malformed input as several unrelated types
        raise DocumentExtractionError(f"extract_text: '{label}' could not be read as a Word document: "
                                      f"{type(exc).__name__}: {exc}") from exc


def _pptx_shapes_text(shapes: Iterable[Any]) -> str:
    """Text of a `python-pptx` shape collection — a slide's shapes, or a group's members — in document order.

    Recursive, so text inside a grouped shape is not lost.
    """
    chunks = []
    for shape in shapes:
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
            chunks.append(_pptx_shapes_text(shape.shapes))
        elif shape.has_table:
            for row in shape.table.rows:
                chunks.append(_TABLE_CELL_SEPARATOR.join(cell.text for cell in row.cells))
        elif shape.has_text_frame:
            chunks.append(shape.text_frame.text)
    return "\n".join(chunk for chunk in chunks if chunk)


def _extract_pptx(source: BinaryIO, label: str) -> str:
    try:
        presentation = pptx.Presentation(source)
        chunks = []
        for slide in presentation.slides:
            chunks.append(_pptx_shapes_text(slide.shapes))
            # Presenter notes are included. Slides tend to carry the claim and the notes the argument for it, so
            # on lecture or conference decks the notes are often where the substance actually is — and a reader
            # asking the corpus a question wants that, projected or not.
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
                chunks.append(slide.notes_slide.notes_text_frame.text)
        return "\n".join(chunk for chunk in chunks if chunk)
    except Exception as exc:  # noqa: BLE001 -- python-pptx surfaces malformed input as several unrelated types
        raise DocumentExtractionError(f"extract_text: '{label}' could not be read as a PowerPoint presentation: "
                                      f"{type(exc).__name__}: {exc}") from exc


# ODF marks a paragraph as `text:p` and a heading as `text:h`; both are leaves as far as we are concerned.
_ODF_PARAGRAPH_QNAMES = frozenset([(odf.namespaces.TEXTNS, "p"),
                                   (odf.namespaces.TEXTNS, "h")])


def _odf_paragraphs(node: Any) -> Iterator[Any]:
    """Yield the paragraph and heading elements under an ODF node, in document order.

    Stops descending at each one it yields. `odf.teletype.extractText` already gathers the text of everything
    nested *inside* a paragraph — a footnote body, a text box anchored as a character — so going deeper would
    emit that text a second time.
    """
    for child in getattr(node, "childNodes", ()):
        if getattr(child, "qname", None) in _ODF_PARAGRAPH_QNAMES:
            yield child
        else:
            yield from _odf_paragraphs(child)


def _extract_odf(source: BinaryIO, label: str) -> str:
    # One reader for both OpenDocument formats we accept: a text document and a presentation differ in how the
    # body is structured, but a walk that only looks for paragraphs does not have to care which it is holding.
    try:
        document = odf.opendocument.load(source)
        return "\n".join(odf.teletype.extractText(paragraph) for paragraph in _odf_paragraphs(document.body))
    except Exception as exc:  # noqa: BLE001 -- odfpy surfaces malformed input as several unrelated types
        raise DocumentExtractionError(f"extract_text: '{label}' could not be read as an OpenDocument file: "
                                      f"{type(exc).__name__}: {exc}") from exc


# How much of a page's text content selection must keep before we believe it selected rather than truncated.
#
# Readability extraction picks *one* main content block, which is the right call for a news page wrapped in
# navigation and comments, and the wrong one for a saved document that holds several: a multi-chapter story
# arrives as one `<article>` per chapter and comes back as a single chapter, with nothing to indicate that the
# other twelve were dropped. Silent, and catastrophic for a retrieval index — the document is present, findable
# and 6% complete.
#
# Measured over 32 saved multi-chapter pages: where selection was correct it kept 82-100% of the page's text,
# and where it truncated it kept 2-32%. Nothing landed in between, so the halfway point has better than 1.5x
# margin on both sides. The two failure directions are not symmetric either, which is what justifies erring
# low: falling back needlessly costs a page's navigation chrome, while not falling back costs the document.
_HTML_CONTENT_SELECTION_MIN_RETENTION = 0.5


def _extract_html(source: BinaryIO, label: str) -> str:
    # `trafilatura` is imported here rather than at module level because it costs about 0.3 s — three times the
    # whole office stack put together — and most sessions never open an HTML document. `raven.server.modules
    # .webfetch` defers it at its use site for the same reason, and this is the same extractor doing the same
    # job on bytes that arrived from disk instead of from the network.
    import trafilatura  # noqa: PLC0415 -- deferred: keep the readability stack off the common import path

    # Handed over as bytes, not text: an HTML file declares its own encoding in a meta tag or an XML
    # declaration, and a page saved off the web is about as likely to be Latin-1 as UTF-8. Decoding it ourselves
    # would mean either guessing or raising on a file that is perfectly well-formed and says so in its header.
    raw = source.read()
    try:
        # Markdown rather than bare text, so headings, lists and tables survive as structure the chunker can see
        # — the document database already accepts Markdown as an input format, so nothing downstream is
        # surprised by it. `favor_recall` errs toward keeping borderline content: for a retrieval index, a
        # paragraph wrongly kept costs far less than one wrongly discarded.
        body = trafilatura.extract(raw, output_format="markdown", include_tables=True, favor_recall=True) or ""

        # Guard against content selection having discarded content rather than chrome. `html2txt` takes the
        # whole page instead of choosing a block, so it is the denominator that says how much went missing.
        whole_page = trafilatura.html2txt(raw) or ""
        if len(body) < _HTML_CONTENT_SELECTION_MIN_RETENTION * len(whole_page):
            # Re-extract per `<article>`, which is the remedy that keeps the Markdown: selection went wrong
            # by choosing among the page's compositions, so hand it each one on its own and let it choose
            # nothing. Whole-page text is the last resort, for a page whose blocks are not marked up as
            # articles — flat and complete beats shapely and truncated, but it is the worse of the two.
            per_article = _extract_html_articles(trafilatura, raw)
            recovered = per_article if len(per_article) >= len(body) else ""
            logger.info(f"_extract_html: '{label}': content selection kept {len(body)} of {len(whole_page)} "
                        f"characters, which reads as truncation rather than boilerplate removal; "
                        f"{'re-extracting per article' if recovered else 'falling back to whole-page text'}.")
            body = recovered or whole_page

        title = _html_title(trafilatura, raw)
    except Exception as exc:  # noqa: BLE001 -- trafilatura surfaces malformed input as several unrelated types
        raise DocumentExtractionError(f"extract_text: '{label}' could not be read as an HTML document: "
                                      f"{type(exc).__name__}: {exc}") from exc

    # Readability extraction keeps the article's own headings but drops `<title>`, which on a saved page is
    # often the only thing that names it — the filename frequently does not. Skipped when the body already
    # opens with that same heading, which is the common case for a page whose `<h1>` restates its title.
    if title and not _body_opens_with(body, title):
        return f"# {title}\n\n{body}" if body else f"# {title}"
    return body


def _extract_html_articles(trafilatura: Any, raw: bytes) -> str:
    """Extract each `<article>` on a page separately, as Markdown, and join them in document order.

    For a page that holds several compositions rather than one — a serial archived as one article per
    chapter, a digest, a thread — extracting the page as a whole makes the readability pass choose between
    them, and it discards everything it did not choose. Handing it one article at a time removes the choice.

    Returns `""` when the page has no `<article>` elements, or when none of them yields any text, so a caller
    can treat an empty result as "this page is not built that way" with no special case.
    """
    import lxml.html  # noqa: PLC0415 -- deferred alongside trafilatura, which is what brings lxml in

    try:
        tree = lxml.html.fromstring(raw)
    except Exception:  # noqa: BLE001 -- lxml surfaces malformed markup as several unrelated types
        return ""

    parts = []
    for article in tree.findall(".//article"):
        body = trafilatura.extract(lxml.html.tostring(article, encoding="unicode"),
                                   output_format="markdown", include_tables=True, favor_recall=True) or ""
        if not body:
            continue
        # The readability pass treats an article's own heading as its title and drops it, the same way it
        # drops the page's `<title>` — so recover it here for the same reason, one level below the page
        # heading the caller prepends. On a serial these are the chapter titles, and they are worth having
        # in a retrieval index: "which chapter is the one about the shutdown" is a question people ask.
        heading = _first_heading(article)
        if heading and not _body_opens_with(body, heading):
            body = f"## {heading}\n\n{body}"
        parts.append(body)
    return "\n\n".join(parts)


def _first_heading(element: Any) -> str | None:
    """Text of the first `<h1>`–`<h6>` anywhere inside an lxml element, or `None` if it has no heading."""
    for node in element.iter("h1", "h2", "h3", "h4", "h5", "h6"):
        text = " ".join(node.text_content().split())
        if text:
            return text
    return None


def _body_opens_with(body: str, title: str) -> bool:
    """Whether `body` already starts with `title`, so prepending it as a heading would state it twice.

    Compares against the start of the first line rather than the whole of it, because the extractor does not
    always give the title a line of its own. A short page can come back as one flat run of text whose first
    line opens with the title and continues into the first paragraph; asking for equality then answers "no"
    on a body that visibly does begin with the title. Prefix matching covers both shapes.

    The asymmetry is deliberate. A false "yes" costs a missing heading on a page that names itself in its
    opening words anyway; a false "no" prints the title twice, which is the defect this guards.
    """
    return body.lstrip().split("\n", 1)[0].lstrip("#").strip().startswith(title)


def _html_title(trafilatura: Any, raw: bytes) -> str | None:
    """Best-effort page title from HTML bytes. `None` when there isn't one, or the metadata pass fails."""
    try:
        metadata = trafilatura.extract_metadata(raw)
    except Exception:  # noqa: BLE001 -- a title is a nicety; never let its absence cost us the body text
        return None
    title = getattr(metadata, "title", None) if metadata is not None else None
    return title.strip() if title else None


# The formats that need a real parser, and who parses them. Anything not listed here is read as plain text, so
# this table plus `_PLAINTEXT_EXTS` is the whole capability list — `supported_extensions` derives from the two
# rather than repeating them, which is what keeps the dispatch and the advertised formats from drifting apart.
_EXTRACTORS = {".pdf": _extract_pdf,
               ".docx": _extract_docx,
               ".pptx": _extract_pptx,
               ".odt": _extract_odf,
               ".odp": _extract_odf,
               ".html": _extract_html,
               ".htm": _extract_html}


def supported_extensions() -> tuple[str, ...]:
    """Return all file extensions `extract_text` can handle (lowercase, with the leading dot)."""
    return _PLAINTEXT_EXTS + tuple(_EXTRACTORS)


def is_supported(path: str | pathlib.Path) -> bool:
    """Return whether `extract_text` recognizes `path`'s file extension."""
    return pathlib.Path(path).suffix.lower() in supported_extensions()


def repair_surrogates(text: str) -> str:
    """Make `text` UTF-8 encodable, recovering what can be recovered.

    A Python `str` holds code points, so any U+D800–U+DFFF in one is a UTF-16 artifact that leaked in
    undecoded — and `str.encode("utf-8")` refuses it. That refusal is the whole problem: extraction
    succeeds, the text looks fine in a REPL, and the failure surfaces much later at whatever first tries
    to put the text on a wire or in a file.

    Two cases, and they want opposite treatment:

      - A surrogate *pair* held as two code points is the real character with its UTF-16 encoding
        showing. Re-reading the string as the UTF-16 code units it is recovers it — U+D835 U+DC34
        becomes 𝐴. Worth doing rather than deleting: mathematical alphanumerics are the common source,
        and dropping them silently unmaths a scientific paper.
      - An *unpaired* surrogate encodes nothing and has to go.

    Both fall out of one round trip, so mixed text needs no special case. Well-formed text is returned
    unchanged, astral characters included.
    """
    try:
        text.encode("utf-8")
        return text
    except UnicodeEncodeError:
        return text.encode("utf-16", "surrogatepass").decode("utf-16", "ignore")


def _extract(source: BinaryIO, name: str, label: str) -> str | None:
    """Dispatch on `name`'s extension, extract, and apply the tail both entry points share.

    `name` decides the format and `label` names the source in messages. They are the same string for a
    file and differ for anything else — a downloaded document whose URL is worth reporting, say.
    """
    # Plain text for the known text extensions, and as a defensive fallback for anything else (the ingester
    # filters by extension upstream, so an unknown suffix reaching here is already an unusual case).
    extract = _EXTRACTORS.get(pathlib.Path(name).suffix.lower(), _extract_plaintext)
    text = extract(source, label).strip()
    # Every extractor funnels through here, which is the point: a backend that hands back undecoded UTF-16
    # is a property of that backend, not of the caller, and every consumer downstream assumes valid text.
    repaired = repair_surrogates(text)
    if repaired != text:
        logger.warning(f"extract_text: '{label}' extracted with UTF-16 surrogates in the text; repaired.")
    # The same argument one step further. An extracted document is somebody else's text, and it goes on to
    # be embedded, indexed and put in front of a model — so a zero-width space or a directional override in
    # it is a property of the document, and every consumer downstream assumes what it gets is what a person
    # would see. The web paths already normalize for exactly this reason; a PDF somebody dropped into a
    # corpus is the same kind of input as a page somebody fetched.
    #
    # After the surrogate repair, so it works on text that is already well-formed.
    normalized = common_text.normalize(repaired)
    if normalized != repaired:
        logger.info(f"extract_text: '{label}' carried invisible or control characters; normalized.")
    return normalized or None


def extract_text(path: str | pathlib.Path) -> str | None:
    """Extract indexable plaintext from a document file. See the module docstring for the full contract."""
    p = pathlib.Path(path).expanduser()
    if not p.exists():
        raise FileNotFoundError(f"extract_text: no such file: '{p}'")
    with open(p, "rb") as source:
        return _extract(source, name=p.name, label=str(p))


def extract_text_from_bytes(raw: bytes, name: str) -> str | None:
    """Extract indexable plaintext from a document held in memory. Same contract, minus `FileNotFoundError`.

    `raw`: The document's bytes.

    `name`: What the document is called — the original filename, or anything else carrying the right
            extension. **This is what selects the reader**, since bytes do not announce their format, and
            it is also what names the document in any error. A name with no recognized extension is read as
            plain text, as an unrecognized file is.

    For a document that never becomes a file: an attachment held in a chat, a fetch that came off the
    network, a page rendered in memory. Writing it to a temp file first would work — every reader here
    happily takes a path — but that is a filesystem round trip to satisfy an interface rather than a need,
    and one that has to be cleaned up afterwards on every path including the failing ones.
    """
    return _extract(io.BytesIO(raw), name=name, label=name)


class Extractor:
    """A document reader together with the file formats it can read.

    These are two halves of one fact, and a consumer that takes them as separate arguments lets them drift:
    a reader handed a format it cannot parse yields mojibake or an exception, and a format list that outruns
    its reader is how a PDF gets indexed as line noise. Neither mistake announces itself — the document is
    in the index, findable, and wrong. Passing the pair as one object means a caller cannot widen one
    without the other.

    Instances are callable, so one can be used anywhere a bare read function was.

    Two ready-made ones live in this module: `PLAINTEXT`, which admits only the formats readable without a
    parser, and `ALL_FORMATS`, which admits everything `extract_text` handles. Narrow either with
    `restricted_to`; there is deliberately no way to widen one, since widening is the direction that breaks.
    """

    def __init__(self, read: Callable[[str | pathlib.Path], str | None], extensions: Iterable[str]) -> None:
        """`read`: Takes a path, returns its plaintext, or `None` if the file holds no extractable text.

        `extensions`: The file extensions `read` can handle, lowercase, with the leading dot.
        """
        self.read = read
        self.extensions = tuple(sorted({ext.lower() for ext in extensions}))

    def handles(self, path: str | pathlib.Path) -> bool:
        """Whether this extractor recognizes `path`'s file extension. Says nothing about the file's contents."""
        return pathlib.Path(path).suffix.lower() in self.extensions

    def restricted_to(self, extensions: Iterable[str]) -> "Extractor":
        """Return a copy of this extractor admitting only `extensions`, same reader.

        For a caller that reads fewer formats than it could — a user who would rather their document folder
        ignored PDFs, say. Anything asked for that this extractor cannot read is dropped with a warning
        rather than honored, which is the point: the result can only ever be a subset of what the reader
        actually handles, so the pair cannot be made inconsistent from the outside.
        """
        wanted = {ext.lower() for ext in extensions}
        unreadable = wanted - set(self.extensions)
        if unreadable:
            logger.warning(f"Extractor.restricted_to: ignoring {sorted(unreadable)}: no reader for "
                           f"{'them' if len(unreadable) > 1 else 'it'} here. Readable formats are "
                           f"{list(self.extensions)}.")
        return Extractor(read=self.read, extensions=wanted & set(self.extensions))

    def __call__(self, path: str | pathlib.Path) -> str | None:
        return self.read(path)

    def __repr__(self) -> str:
        return f"<{type(self).__name__}: {' '.join(self.extensions)}>"


# The reader is the same dispatcher in both; what differs is which formats each admits. So the plain-text one
# is not a lesser reader, it is the same reader told to stay out of the formats that need a parser installed.
PLAINTEXT = Extractor(read=extract_text, extensions=_PLAINTEXT_EXTS)
ALL_FORMATS = Extractor(read=extract_text, extensions=supported_extensions())
