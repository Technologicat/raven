"""Shared test utilities for Raven."""

__all__ = ["approx",

           "make_minimal_pdf",
           "make_textless_pdf",
           "make_docx",
           "make_pptx",
           "make_odt",
           "make_odp"]

import io
from typing import Callable

import docx
import odf.draw
import odf.opendocument
import odf.style
import odf.text
import pptx
import pptx.util


def approx(a, b, tol=0.01):
    """Check approximate float equality."""
    return abs(a - b) < tol


def make_minimal_pdf(text: str) -> bytes:
    """Build a minimal but valid single-page PDF whose text layer is `text`, with a correct xref table.

    `text` is embedded in a `Tj` text-showing operator, so a born-digital PDF text extractor (pypdf) reads it
    back verbatim. Restricted to Latin-1 for the content stream. Returns the PDF file bytes, ready to write to
    disk — used as an on-the-fly fixture so tests need no committed binary.
    """
    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R "
        b"/Resources << /Font << /F1 5 0 R >> >> >>",
        None,  # contents stream, filled in below (its length depends on `text`)
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    stream = b"BT /F1 24 Tf 72 700 Td (" + text.encode("latin-1") + b") Tj ET"
    objs[3] = b"<< /Length %d >>\nstream\n%s\nendstream" % (len(stream), stream)
    return _assemble_pdf(objs)


def make_textless_pdf() -> bytes:
    """Build a valid single-page PDF with no text-showing operators — a stand-in for a scanned/image-only page.

    A born-digital text extractor finds no text in it, so it exercises the "parses cleanly but yields no text"
    path (which should come back as empty rather than as an error).
    """
    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R >>",
        b"<< /Length 0 >>\nstream\n\nendstream",  # empty content stream: nothing to extract
    ]
    return _assemble_pdf(objs)


def make_docx(blocks: list) -> bytes:
    """Build a minimal `.docx` document. Returns the file bytes, ready to write to disk.

    Each block is either a string (one paragraph) or a list of rows, each row a list of cell strings (one
    table). Mixing the two exercises document order, which is the part of Word extraction most easily got
    wrong — a table can only be read back in its right place if the body is walked rather than queried.
    """
    document = docx.Document()
    for block in blocks:
        if isinstance(block, str):
            document.add_paragraph(block)
        else:
            table = document.add_table(rows=len(block), cols=len(block[0]))
            for row_index, row in enumerate(block):
                for column_index, cell_text in enumerate(row):
                    table.cell(row_index, column_index).text = cell_text
    return _to_bytes(document.save)


def make_pptx(slides: list[dict]) -> bytes:
    """Build a minimal `.pptx` presentation. Returns the file bytes, ready to write to disk.

    Each slide is a dict, and every key is optional:

      - `"text"`: list of strings, each placed in its own text box;
      - `"table"`: list of rows, each row a list of cell strings;
      - `"group"`: list of strings, placed in text boxes inside one grouped shape (so that recursion into
        groups is exercised — text hidden in a group is otherwise silently dropped);
      - `"notes"`: string, the slide's presenter notes.
    """
    presentation = pptx.Presentation()
    blank_layout = presentation.slide_layouts[6]
    for slide_spec in slides:
        slide = presentation.slides.add_slide(blank_layout)
        # Stack the shapes down the slide. Nothing here reads geometry back, but non-overlapping boxes keep a
        # generated fixture openable in a real editor when a test result needs to be looked at by eye.
        top = pptx.util.Inches(0.5)
        row_height = pptx.util.Inches(1.0)
        for text in slide_spec.get("text", []):
            textbox = slide.shapes.add_textbox(pptx.util.Inches(0.5), top, pptx.util.Inches(6), row_height)
            textbox.text_frame.text = text
            top += row_height
        if (rows := slide_spec.get("table", None)) is not None:
            graphic_frame = slide.shapes.add_table(len(rows), len(rows[0]),
                                                   pptx.util.Inches(0.5), top, pptx.util.Inches(6), row_height)
            for row_index, row in enumerate(rows):
                for column_index, cell_text in enumerate(row):
                    graphic_frame.table.cell(row_index, column_index).text = cell_text
            top += row_height
        if (grouped_texts := slide_spec.get("group", None)) is not None:
            group = slide.shapes.add_group_shape()
            for text in grouped_texts:
                textbox = group.shapes.add_textbox(pptx.util.Inches(0.5), top, pptx.util.Inches(6), row_height)
                textbox.text_frame.text = text
                top += row_height
        if (notes := slide_spec.get("notes", None)) is not None:
            slide.notes_slide.notes_text_frame.text = notes
    return _to_bytes(presentation.save)


def make_odt(blocks: list[tuple[str, str]]) -> bytes:
    """Build a minimal `.odt` text document. Returns the file bytes, ready to write to disk.

    Each block is `("p", text)` for a paragraph or `("h", text)` for a heading. ODF marks the two with different
    element names, so interleaving them checks that both are collected, and collected in order.
    """
    document = odf.opendocument.OpenDocumentText()
    for kind, text in blocks:
        element = odf.text.H(outlinelevel=1, text=text) if kind == "h" else odf.text.P(text=text)
        document.text.addElement(element)
    return _to_bytes(document.write)


def make_odp(slides: list[list[str]]) -> bytes:
    """Build a minimal `.odp` presentation. Returns the file bytes, ready to write to disk.

    Each slide is a list of paragraph strings, placed in one text box on that slide. The text therefore sits
    several levels down (page → frame → text box → paragraph), which is the arrangement a presentation uses and
    a text document does not.
    """
    document = odf.opendocument.OpenDocumentPresentation()
    # A `draw:page` names the master page it is drawn on, so there has to be one for the file to be well-formed
    # — even though nothing about the text extraction looks at it.
    document.masterstyles.addElement(odf.style.MasterPage(name="Fixture", pagelayoutname="FixtureLayout"))
    for paragraphs in slides:
        page = odf.draw.Page(masterpagename="Fixture")
        document.presentation.addElement(page)
        frame = odf.draw.Frame(x="1cm", y="1cm", width="20cm", height="5cm")
        page.addElement(frame)
        textbox = odf.draw.TextBox()
        frame.addElement(textbox)
        for text in paragraphs:
            textbox.addElement(odf.text.P(text=text))
    return _to_bytes(document.write)


def _to_bytes(save: Callable[[io.BytesIO], None]) -> bytes:
    """Call a document object's save/write method against an in-memory buffer, and return the bytes written."""
    buffer = io.BytesIO()
    save(buffer)
    return buffer.getvalue()


def _assemble_pdf(objs: list) -> bytes:
    """Assemble a list of PDF object bodies (1-indexed) into a complete PDF with a correct xref table."""
    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, body in enumerate(objs, start=1):
        offsets.append(len(out))
        out += b"%d 0 obj\n%s\nendobj\n" % (i, body)
    xref_pos = len(out)
    out += b"xref\n0 %d\n" % (len(objs) + 1)
    out += b"0000000000 65535 f \n"
    for off in offsets:
        out += b"%010d 00000 n \n" % off
    out += b"trailer\n<< /Size %d /Root 1 0 R >>\nstartxref\n%d\n%%%%EOF\n" % (len(objs) + 1, xref_pos)
    return bytes(out)
